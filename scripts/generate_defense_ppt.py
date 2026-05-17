from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output"
PPTX_PATH = OUT_DIR / "戴驰峰-结项答辩PPT.pptx"
NOTES_PATH = OUT_DIR / "戴驰峰-结项答辩PPT-讲稿.md"

SLIDE_W = 13.333333
SLIDE_H = 7.5

COLORS = {
    "navy": "12355B",
    "blue": "2D6CDF",
    "teal": "167C80",
    "mint": "EAF7F3",
    "orange": "E76F51",
    "red": "C9413E",
    "green": "2A9D8F",
    "yellow": "F4A261",
    "ink": "17212B",
    "muted": "5E6C7A",
    "line": "D8E0EA",
    "panel": "FFFFFF",
    "bg": "F6F8FB",
    "soft_blue": "EAF1FF",
    "soft_orange": "FFF1E8",
    "soft_red": "FDEDEC",
    "soft_green": "EAF7F3",
}


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def find_image(patterns: Sequence[str]) -> Path:
    for pattern in patterns:
        matches = sorted(ROOT.glob(pattern))
        if matches:
            return matches[0]
    raise FileNotFoundError(f"未找到图片：{patterns}")


IMAGES = {
    "front_overview": find_image(["image/fig-5-1-system-frontend-overview.png"]),
    "hash_verify": find_image(["image/fig-2-2-hash-integrity-verification.png"]),
    "hybrid_storage": find_image(["image/fig-3-1-chain-offchain-hybrid-storage-model.png"]),
    "three_hash": find_image(["image/fig-3-3-three-party-hash-comparison-flow.png"]),
    "architecture": find_image(["image/*系统总体架构图*去标题.png", "doc/architecture-v4-classic.png"]),
    "dashboard": find_image(["doc/screenshots/09_dashboard_with_data.png", "image/fig-5-1-system-frontend-overview.png"]),
    "audit_chart": find_image(["image/fig-5-3-batch-audit-time-comparison.png"]),
    "tamper_loop": find_image(["image/fig-5-7-tamper-detection-loop.png"]),
}


SLIDE_NOTES = [
    (
        "题目与研究目标",
        "各位老师好，我的毕业设计题目是《基于区块链的可信任务日志审计系统设计与实现》。"
        "本课题面向任务执行日志的可信审计场景，目标是把日志采集、链下存储、链上哈希存证、"
        "审计比对、异常告警和前端展示串联成一个完整闭环。"
    ),
    (
        "背景与问题",
        "在任务执行和运维审计中，日志是判断任务是否正常执行的重要依据。但传统中心化日志存在一个问题："
        "日志原文和审计依据往往保存在同一个管理域内，一旦数据库被攻击或高权限人员误用权限，"
        "事后很难证明日志内容是否保持原样。因此本课题关注的是日志的事后可验证性。"
    ),
    (
        "总体方案",
        "系统采用链上链下混合存储。日志原文保存在链下数据库中，便于查询和展示；同时对日志原文计算 SHA-256 哈希，"
        "把哈希摘要和任务标识写入区块链合约。这样既避免原始日志直接上链带来的成本和隐私问题，"
        "又保留了相对独立、难以随意篡改的链上校验依据。"
    ),
    (
        "三方哈希比对机制",
        "系统的核心审计机制是 expectedHash、actualHash 和 onChainHash 三方比对。"
        "expectedHash 来自日志提交阶段保存的链下哈希记录，actualHash 是审计时对当前日志原文重新计算得到的哈希，"
        "onChainHash 来自 LogRegistry 合约。三者一致则审计通过；哈希不一致则判定为 failed 并生成 hash_mismatch 告警；"
        "链上记录暂不可用时标记为 pending。"
    ),
    (
        "系统架构与模块实现",
        "工程实现上，系统包括 Agent、Server、SQLite、LogRegistry 和 Web 五个部分。"
        "Agent 负责自动增量采集日志，Server 负责日志入库、哈希计算、合约调用、审计比对和告警生成，"
        "SQLite 保存链下业务数据，LogRegistry 保存链上哈希存证，Web 前端负责展示日志、审计和告警结果。"
    ),
    (
        "系统运行演示",
        "接下来用一分钟展示系统运行效果。视频中依次展示日志生成与上链存证、日志中心查看、批量审计、"
        "篡改检测以及 hash_mismatch 告警生成。这个过程可以直观看到系统从日志产生到异常反馈的完整闭环。"
    ),
    (
        "功能验证与实验结果",
        "实验结果显示，LogRegistry 合约 8 项测试全部通过；100 条日志连续提交成功率为 100%，平均响应时间约 107.03 毫秒，"
        "吞吐量约 9.33 条每秒。批量审计在 100、500 和 1000 条规模下各执行 5 轮均成功，"
        "篡改检测实验中系统能够将被修改的日志判定为 failed，并自动生成 hash_mismatch 告警。"
    ),
    (
        "总结与展望",
        "总结来看，本系统完成了一个基于区块链的可信任务日志审计原型。主要特点是链下日志原文与链上哈希摘要结合，"
        "并通过三方哈希比对实现日志完整性审计。后续可以接入真实链环境，优化批量审计性能，并完善多角色权限和告警处理流程。"
    ),
    (
        "备用答辩问题",
        "本页用于答辩问答备用。重点准备为什么要使用区块链、为什么不直接上链日志原文、如何判断篡改、"
        "pending 状态含义以及当前系统不足。"
    ),
]


def set_fill(shape, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_line(shape, color: str = "D8E0EA", width: float = 1.0, transparency: int = 0) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = COLORS["ink"],
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Arial",
    valign=MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
) -> object:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0)
    frame.margin_right = Inches(0)
    frame.margin_top = Inches(0)
    frame.margin_bottom = Inches(0)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_multiline(
    slide,
    lines: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 15,
    color: str = COLORS["ink"],
    bullet: bool = False,
    gap_before: float = 4,
) -> object:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0)
    frame.margin_bottom = Inches(0)
    for idx, line in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(gap_before)
        if bullet:
            p.level = 0
            p._p.get_or_add_pPr().insert(0, p._p._new_buChar())
            p._p.pPr.buChar.set("char", "•")
    return box


def add_title(slide, title: str, subtitle: str | None = None, page_no: int | None = None) -> None:
    add_text(slide, title, 0.62, 0.34, 7.8, 0.44, size=23, bold=True, color=COLORS["navy"])
    if subtitle:
        add_text(slide, subtitle, 0.64, 0.82, 7.4, 0.25, size=8.5, color=COLORS["muted"])
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.02), Inches(1.12), Inches(0.05))
    set_fill(accent, COLORS["orange"])
    set_line(accent, COLORS["orange"], transparency=100)
    if page_no is not None:
        add_text(slide, f"{page_no:02d}", 12.18, 0.36, 0.55, 0.2, size=9, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_footer(slide) -> None:
    add_text(slide, "基于区块链的可信任务日志审计系统", 0.62, 7.16, 4.0, 0.18, size=7.5, color="8A97A6")
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(7.05), Inches(12.1), Inches(0.01))
    set_fill(line, "E3E8EF")
    set_line(line, "E3E8EF", transparency=100)


def add_panel(slide, x: float, y: float, w: float, h: float, fill: str = "FFFFFF", line: str = "D8E0EA", radius=True) -> object:
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, line, width=0.8)
    return shape


def add_badge(slide, text: str, x: float, y: float, w: float, fill: str, color: str = "FFFFFF") -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    set_fill(shape, fill)
    set_line(shape, fill, transparency=100)
    add_text(slide, text, x + 0.08, y + 0.08, w - 0.16, 0.16, size=8.2, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_picture_fit(slide, image_path: Path, x: float, y: float, w: float, h: float, border: bool = True) -> None:
    with Image.open(image_path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    pic_w = iw * scale
    pic_h = ih * scale
    px = x + (w - pic_w) / 2
    py = y + (h - pic_h) / 2
    slide.shapes.add_picture(str(image_path), Inches(px), Inches(py), width=Inches(pic_w), height=Inches(pic_h))
    if border:
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        rect.fill.background()
        set_line(rect, "CCD6E2", width=0.8, transparency=5)


def add_metric(slide, value: str, label: str, x: float, y: float, w: float, color: str) -> None:
    add_panel(slide, x, y, w, 0.74, "FFFFFF", "DFE7EF")
    add_text(slide, value, x + 0.12, y + 0.12, w - 0.24, 0.28, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.08, y + 0.47, w - 0.16, 0.16, size=7.7, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def add_flow_chip(slide, text: str, x: float, y: float, w: float, fill: str, color: str = COLORS["navy"]) -> None:
    add_panel(slide, x, y, w, 0.5, fill, "D8E0EA")
    add_text(slide, text, x + 0.08, y + 0.17, w - 0.16, 0.14, size=8.5, bold=True, color=color, align=PP_ALIGN.CENTER)


def add_arrow(slide, x: float, y: float) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(0.34), Inches(0.18))
    set_fill(shape, "AAB6C4")
    set_line(shape, "AAB6C4", transparency=100)


def create_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, COLORS["bg"])
    add_panel(slide, 0.0, 0.0, 13.333, 7.5, COLORS["bg"], COLORS["bg"], radius=False)
    add_text(slide, "基于区块链的可信任务日志审计系统", 0.62, 0.74, 6.95, 1.0, size=31, bold=True, color=COLORS["navy"])
    add_text(slide, "设计与实现", 0.62, 1.78, 3.2, 0.45, size=31, bold=True, color=COLORS["orange"])
    add_text(slide, "结项答辩 / 5分钟汇报", 0.65, 2.42, 2.5, 0.25, size=11, color=COLORS["muted"])
    add_panel(slide, 0.62, 3.05, 5.55, 1.15, "FFFFFF", "D8E0EA")
    add_text(slide, "研究目标", 0.9, 3.27, 1.0, 0.22, size=10, color=COLORS["teal"], bold=True)
    add_text(slide, "构建采集、存证、审计、告警、展示一体化的可信任务日志审计原型系统", 0.9, 3.58, 4.8, 0.34, size=16, color=COLORS["ink"], bold=True)
    add_text(slide, "学生：戴驰峰    课题方向：区块链存证 / 日志审计 / 智能合约", 0.65, 6.65, 6.5, 0.25, size=10.2, color=COLORS["muted"])
    add_panel(slide, 7.35, 0.62, 5.25, 5.85, "FFFFFF", "D8E0EA")
    add_picture_fit(slide, IMAGES["front_overview"], 7.56, 0.84, 4.83, 5.38)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, COLORS["bg"])
    add_title(slide, "中心化日志审计的可信性问题", "日志能保存，不代表日志可被独立验证", 2)
    add_footer(slide)
    add_text(slide, "关键问题", 0.72, 1.42, 1.3, 0.22, size=10, color=COLORS["orange"], bold=True)
    add_text(slide, "事后如何证明日志内容是否与存证时一致？", 0.72, 1.78, 5.25, 0.58, size=26, color=COLORS["navy"], bold=True)
    cards = [
        ("存储集中", "日志原文与审计依据常在同一数据库或文件系统中", COLORS["soft_blue"], COLORS["blue"]),
        ("易被修改", "数据库被攻击或高权限人员误用时，日志可能被修改、删除或伪造", COLORS["soft_red"], COLORS["red"]),
        ("缺少证明", "仅依赖数据库自身，难以提供独立可信的完整性证据", COLORS["soft_orange"], COLORS["orange"]),
    ]
    for i, (head, body, fill, color) in enumerate(cards):
        y = 2.78 + i * 1.12
        add_panel(slide, 0.72, y, 5.45, 0.84, fill, "DDE5EF")
        add_text(slide, head, 0.98, y + 0.14, 1.2, 0.18, size=12, color=color, bold=True)
        add_text(slide, body, 2.02, y + 0.13, 3.75, 0.35, size=12, color=COLORS["ink"])
    add_panel(slide, 6.74, 1.32, 5.92, 4.96, "FFFFFF", "D8E0EA")
    add_picture_fit(slide, IMAGES["hash_verify"], 7.02, 1.66, 5.35, 3.0)
    add_text(slide, "解决思路", 7.05, 4.96, 0.95, 0.2, size=10, color=COLORS["teal"], bold=True)
    add_text(slide, "链下保存日志原文，链上保存哈希摘要，让日志完整性具备可验证依据。", 7.05, 5.26, 4.95, 0.4, size=14, color=COLORS["ink"], bold=True)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, COLORS["bg"])
    add_title(slide, "链下原文 + 链上哈希的混合存证方案", "兼顾查询效率、隐私保护、存储成本与审计可信性", 3)
    add_footer(slide)
    add_panel(slide, 0.68, 1.34, 7.1, 4.95, "FFFFFF", "D8E0EA")
    add_picture_fit(slide, IMAGES["hybrid_storage"], 0.93, 1.6, 6.58, 3.7)
    add_flow_chip(slide, "日志文件", 0.95, 5.58, 1.03, COLORS["soft_blue"])
    add_arrow(slide, 2.1, 5.74)
    add_flow_chip(slide, "Agent", 2.55, 5.58, 0.92, COLORS["soft_green"], COLORS["teal"])
    add_arrow(slide, 3.62, 5.74)
    add_flow_chip(slide, "Server", 4.04, 5.58, 0.98, COLORS["soft_orange"], COLORS["orange"])
    add_arrow(slide, 5.18, 5.74)
    add_flow_chip(slide, "SQLite + LogRegistry", 5.58, 5.58, 1.75, "F0F3F7")
    add_text(slide, "方案要点", 8.32, 1.45, 1.3, 0.22, size=10, color=COLORS["teal"], bold=True)
    steps = [
        "链下 SQLite 保存日志原文、任务编号、来源路径和采集时间",
        "Server 对日志正文计算 SHA-256 哈希",
        "LogRegistry 合约保存 taskId、logHash、createdAt、submitter",
        "审计时重算哈希，并与链下记录和链上记录比对",
        "异常时写入审计记录并生成 hash_mismatch 告警",
    ]
    add_multiline(slide, steps, 8.32, 1.9, 4.05, 3.25, size=13, color=COLORS["ink"], bullet=False, gap_before=10)
    add_badge(slide, "不直接上链原文", 8.28, 5.56, 1.55, COLORS["red"])
    add_badge(slide, "只上链哈希证据", 10.02, 5.56, 1.68, COLORS["teal"])

    # Slide 4
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, COLORS["bg"])
    add_title(slide, "三方哈希比对审计机制", "expectedHash、actualHash、onChainHash 共同给出审计结论", 4)
    add_footer(slide)
    hashes = [
        ("expectedHash", "提交阶段保存在 log_hash_records 中的日志哈希", COLORS["soft_blue"], COLORS["blue"]),
        ("actualHash", "审计阶段对当前日志原文重新计算得到的哈希", COLORS["soft_green"], COLORS["teal"]),
        ("onChainHash", "从 LogRegistry 合约读取的链上哈希摘要", COLORS["soft_orange"], COLORS["orange"]),
    ]
    for i, (head, body, fill, color) in enumerate(hashes):
        x = 0.72 + i * 2.18
        add_panel(slide, x, 1.38, 1.95, 1.22, fill, "D8E0EA")
        add_text(slide, head, x + 0.14, 1.6, 1.66, 0.18, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.17, 1.92, 1.6, 0.34, size=8.6, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_panel(slide, 0.72, 2.98, 6.28, 2.62, "FFFFFF", "D8E0EA")
    result_rows = [
        ("三者一致", "passed", COLORS["green"], "日志内容与存证时一致，审计通过"),
        ("哈希不一致", "failed", COLORS["red"], "判定为异常，并生成 hash_mismatch 告警"),
        ("链上记录缺失", "pending", COLORS["orange"], "记录为待审计，避免误判为通过"),
    ]
    for i, (cond, status, color, desc) in enumerate(result_rows):
        y = 3.28 + i * 0.68
        add_text(slide, cond, 1.0, y, 1.1, 0.18, size=11.5, color=COLORS["ink"], bold=True)
        add_badge(slide, status, 2.22, y - 0.08, 0.86, color)
        add_text(slide, desc, 3.28, y, 3.25, 0.18, size=11, color=COLORS["muted"])
    add_panel(slide, 7.35, 1.34, 5.25, 4.95, "FFFFFF", "D8E0EA")
    add_picture_fit(slide, IMAGES["three_hash"], 7.57, 1.62, 4.82, 4.36)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, COLORS["bg"])
    add_title(slide, "系统架构与模块实现", "Agent、Server、SQLite、LogRegistry、Web 协同形成可信审计闭环", 5)
    add_footer(slide)
    add_panel(slide, 0.65, 1.26, 7.35, 4.8, "FFFFFF", "D8E0EA")
    add_picture_fit(slide, IMAGES["architecture"], 0.9, 1.52, 6.85, 4.25)
    modules = [
        ("Agent", "增量读取、偏移量持久化、失败重试"),
        ("Server", "日志入库、哈希计算、合约调用、审计告警"),
        ("SQLite", "日志、存证记录、审计记录、告警、Agent状态"),
        ("LogRegistry", "链上哈希存证，AccessControl 控制写入权限"),
        ("Web", "总览、日志中心、审计管理、异常告警展示"),
    ]
    for i, (head, body) in enumerate(modules):
        y = 1.33 + i * 0.92
        add_panel(slide, 8.42, y, 3.9, 0.66, "FFFFFF", "D8E0EA")
        add_text(slide, head, 8.63, y + 0.16, 0.95, 0.16, size=11, bold=True, color=[COLORS["blue"], COLORS["teal"], COLORS["orange"], COLORS["red"], COLORS["navy"]][i])
        add_text(slide, body, 9.68, y + 0.15, 2.35, 0.18, size=9.4, color=COLORS["muted"])
    add_text(slide, "核心接口", 8.48, 6.06, 0.9, 0.16, size=8.7, bold=True, color=COLORS["teal"])
    add_text(slide, "POST /api/logs    POST /api/audits/run    GET /api/alerts", 9.22, 6.06, 3.2, 0.16, size=8.5, color=COLORS["muted"])

    # Slide 6
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, COLORS["navy"])
    slide.shapes.add_picture(str(IMAGES["dashboard"]), Inches(6.6), Inches(0), height=Inches(7.5))
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    set_fill(overlay, "12355B", transparency=12)
    set_line(overlay, "12355B", transparency=100)
    add_text(slide, "系统运行演示", 0.78, 0.82, 4.8, 0.56, size=32, bold=True, color="FFFFFF")
    add_text(slide, "建议在此页插入约 1 分钟演示视频", 0.82, 1.58, 4.1, 0.24, size=12, color="D6E3F3")
    demo_steps = ["日志生成与上链存证", "批量审计与状态展示", "篡改检测与告警生成"]
    for i, step in enumerate(demo_steps):
        y = 2.36 + i * 0.78
        add_panel(slide, 0.82, y, 4.5, 0.52, "FFFFFF", "FFFFFF")
        add_text(slide, f"{i + 1}. {step}", 1.08, y + 0.16, 3.8, 0.15, size=12.5, color=COLORS["navy"], bold=True)
    play = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(5.78), Inches(3.18), Inches(0.72), Inches(0.72))
    play.rotation = 90
    set_fill(play, COLORS["orange"])
    set_line(play, COLORS["orange"], transparency=100)
    add_text(slide, "播放视频", 4.95, 4.05, 1.95, 0.24, size=13, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, COLORS["bg"])
    add_title(slide, "功能验证与实验结果", "验证系统可运行性、稳定性与篡改检测能力", 7)
    add_footer(slide)
    add_metric(slide, "100%", "100条日志提交成功率", 0.76, 1.27, 1.85, COLORS["green"])
    add_metric(slide, "107.03 ms", "平均响应时间", 2.8, 1.27, 1.85, COLORS["blue"])
    add_metric(slide, "9.33 条/秒", "吞吐量", 4.84, 1.27, 1.85, COLORS["orange"])
    add_metric(slide, "8项", "合约测试全部通过", 6.88, 1.27, 1.85, COLORS["teal"])
    add_panel(slide, 0.76, 2.34, 5.12, 3.68, "FFFFFF", "D8E0EA")
    add_text(slide, "实验结论", 1.02, 2.62, 1.1, 0.2, size=11, color=COLORS["teal"], bold=True)
    exp_lines = [
        "日志批量提交：100次成功，失败0次",
        "批量审计：100 / 500 / 1000条，5轮均成功",
        "平均审计耗时：3067.77 / 15659.44 / 35032.13 ms",
        "篡改检测：auditStatus = failed",
        "告警结果：alertGenerated = true，类型为 hash_mismatch",
    ]
    add_multiline(slide, exp_lines, 1.02, 3.0, 4.58, 2.26, size=11.2, color=COLORS["ink"], gap_before=8)
    add_panel(slide, 6.32, 2.34, 5.98, 3.68, "FFFFFF", "D8E0EA")
    add_picture_fit(slide, IMAGES["audit_chart"], 6.58, 2.62, 5.46, 3.12)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, COLORS["bg"])
    add_title(slide, "总结与展望", "完成可信任务日志审计原型，后续面向真实环境继续优化", 8)
    add_footer(slide)
    cols = [
        (
            "创新点",
            COLORS["teal"],
            COLORS["soft_green"],
            [
                "链下日志原文 + 链上哈希摘要的轻量存证模型",
                "expectedHash / actualHash / onChainHash 三方比对",
                "历史合约地址回溯与合约代码存在性校验",
                "采集、存证、审计、告警、展示完整闭环",
            ],
        ),
        (
            "存在不足",
            COLORS["orange"],
            COLORS["soft_orange"],
            [
                "主要基于本地 Hardhat 链，未验证真实链部署成本",
                "SQLite 更适合原型系统，高并发能力有限",
                "批量审计以顺序处理为主，大规模场景仍需优化",
                "当前重点检测内容篡改，未深入分析日志语义风险",
            ],
        ),
        (
            "后续优化",
            COLORS["blue"],
            COLORS["soft_blue"],
            [
                "接入联盟链或公开测试网",
                "引入异步审计队列、并发链上查询和索引优化",
                "完善多角色权限体系和告警处理流程",
                "结合语义分析与风险分级增强异常识别能力",
            ],
        ),
    ]
    for i, (head, color, fill, lines) in enumerate(cols):
        x = 0.72 + i * 4.06
        add_panel(slide, x, 1.32, 3.55, 4.86, "FFFFFF", "D8E0EA")
        add_panel(slide, x + 0.18, 1.58, 3.18, 0.56, fill, "D8E0EA")
        add_text(slide, head, x + 0.38, 1.76, 2.65, 0.17, size=12.2, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_multiline(slide, lines, x + 0.34, 2.44, 2.92, 2.7, size=10.5, color=COLORS["ink"], gap_before=9)
    add_text(slide, "结论：系统验证了链上链下混合存储、三方哈希比对和异常告警机制的可行性。", 0.8, 6.44, 11.6, 0.22, size=12.3, color=COLORS["navy"], bold=True, align=PP_ALIGN.CENTER)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    set_fill(slide.background, COLORS["bg"])
    add_title(slide, "备用答辩问题准备", "可放在最后，不主动讲，老师提问时快速切换", 9)
    add_footer(slide)
    qa = [
        ("为什么要用区块链？", "普通数据库便于查询，但不能提供独立可信的审计依据；区块链在本系统中作为哈希证据源。"),
        ("为什么不直接上链原文？", "日志原文体积大，可能包含路径、错误信息等敏感内容；上链哈希兼顾成本、隐私和可验证性。"),
        ("如何判断日志被篡改？", "审计阶段重新计算 actualHash，并与 expectedHash 和 onChainHash 做三方比对。"),
        ("pending 是什么情况？", "链上记录缺失、合约不可用或环境变化时，系统不直接判定通过，而是保留待审计状态。"),
        ("当前最大不足是什么？", "本地链和 SQLite 更适合原型验证，真实部署还需要链环境、数据库和审计性能优化。"),
    ]
    for i, (q, a) in enumerate(qa):
        y = 1.34 + i * 1.04
        add_panel(slide, 0.86, y, 11.48, 0.78, "FFFFFF", "D8E0EA")
        add_text(slide, q, 1.13, y + 0.16, 2.25, 0.18, size=11.5, bold=True, color=[COLORS["blue"], COLORS["teal"], COLORS["orange"], COLORS["red"], COLORS["navy"]][i])
        add_text(slide, a, 3.56, y + 0.14, 8.2, 0.28, size=10.3, color=COLORS["ink"])

    return prs


def write_notes() -> None:
    lines = [
        "# 戴驰峰-结项答辩PPT讲稿",
        "",
        "总时长建议：约 5 分钟；其中第 6 页播放约 1 分钟系统演示视频。",
        "",
    ]
    for idx, (title, note) in enumerate(SLIDE_NOTES, 1):
        lines.extend([f"## 第 {idx} 页：{title}", "", note, ""])
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = create_deck()
    prs.save(PPTX_PATH)
    write_notes()
    print(PPTX_PATH)
    print(NOTES_PATH)


if __name__ == "__main__":
    main()
